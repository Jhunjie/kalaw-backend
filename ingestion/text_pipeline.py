# backend/ingestion/text_pipeline.py
import uuid, fitz
from pptx import Presentation
from docx import Document
from core.llm import embed

CHUNK_SIZE    = 400   # words per chunk
CHUNK_OVERLAP = 80    # overlapping words between chunks

def extract_text(path: str, ext: str) -> str:
    if ext == ".pdf":
        doc = fitz.open(path)
        return "\n".join(p.get_text() for p in doc)
    elif ext in (".pptx", ".ppt"):
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)
    elif ext in (".docx", ".doc"):
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    return ""

def chunk_text(text: str) -> list[str]:
    words  = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + CHUNK_SIZE])
        if chunk.strip():
            chunks.append(chunk)
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks

async def ingest_document_text(
    file_path, document_id, document_name,
    uploader_name, program, subject
) -> list[dict]:
    from pathlib import Path
    ext  = Path(file_path).suffix.lower()
    text = extract_text(file_path, ext)
    if not text.strip():
        return []

    chunks = chunk_text(text)
    rows   = []
    for i, chunk in enumerate(chunks):
        embed_text = (
            f"{document_name} by {uploader_name} "
            f"({program} - {subject}): {chunk}"
        )
        embedding = await embed(embed_text)
        rows.append({
            "id":            str(uuid.uuid4()),
            "document_id":   document_id,
            "document_name": document_name,
            "uploader_name": uploader_name,
            "program":       program,
            "subject":       subject,
            "chunk_text":    chunk,
            "chunk_index":   i,
            "embedding":     embedding,
        })
    return rows
