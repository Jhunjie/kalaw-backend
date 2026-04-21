"""
CPSU AI - Image Ingestion Pipeline
Extracts images from PDF, PPTX, DOCX files and generates
text descriptions using Llama 3.2 Vision via Ollama.
"""

import os
import io
import base64
import json
import uuid
import httpx
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import Optional

# Document parsing
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from PIL import Image


# ── Config ────────────────────────────────────────────────────────────────────

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
VISION_MODEL    = os.getenv("VISION_MODEL", "llama3.2-vision:11b")
EMBED_MODEL     = os.getenv("EMBED_MODEL",  "nomic-embed-text")
MIN_IMAGE_SIZE  = 100   # px — skip tiny icons/bullets
MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5MB cap before resize


# ── Data models ───────────────────────────────────────────────────────────────

@dataclass
class ImageChunk:
    """One extracted image with its metadata and AI-generated description."""
    chunk_id:        str
    document_id:     str
    document_name:   str
    uploader_name:   str
    program:         str          # e.g. "HM", "IT", "Education"
    subject:         str          # e.g. "Culinary Arts 101"
    source_type:     str          # "pdf" | "pptx" | "docx"
    page_or_slide:   int
    image_index:     int          # position within page/slide
    description:     str          # LLM-generated description
    image_b64:       str          # base64 for storage / display
    width:           int
    height:          int
    embedding:       Optional[list[float]] = None


# ── Ollama helpers ─────────────────────────────────────────────────────────────

async def describe_image(image_b64: str, context_hint: str = "") -> str:
    """
    Send an image to Llama 3.2 Vision and get a rich description.
    context_hint can be the surrounding text from the document.
    """
    system_prompt = (
        "You are an expert academic content analyzer for a Philippine university. "
        "Describe the image in detail focusing on: what the object/content IS, "
        "its academic/educational significance, any labels or text visible, "
        "and its likely use in an educational context. "
        "If it is a utensil, tool, equipment, diagram, chart, or food item — "
        "be very specific about the name, type, and function. "
        "Respond in English but note if the content relates to Filipino/local context."
    )

    user_content = "Describe this image in detail for an educational AI system."
    if context_hint:
        user_content += f" Context from the document: {context_hint[:300]}"

    payload = {
        "model": VISION_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": user_content,
                "images": [image_b64],
            },
        ],
        "stream": False,
        "options": {"temperature": 0.1},
    }

    async with httpx.AsyncClient(timeout=100) as client:
        resp = await client.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            json=payload,
        )
        resp.raise_for_status()
        data = resp.json()
        return data["message"]["content"].strip()


async def embed_text(text: str) -> list[float]:
    """Generate a text embedding using nomic-embed-text via Ollama."""
    async with httpx.AsyncClient(timeout=30.0) as client:
        resp = await client.post(
            f"{OLLAMA_BASE_URL}/api/embeddings",
            json={"model": EMBED_MODEL, "prompt": text},
        )
        resp.raise_for_status()
        return resp.json()["embedding"]


# ── Image utilities ────────────────────────────────────────────────────────────

def pil_to_b64(img: Image.Image, max_bytes: int = MAX_IMAGE_BYTES) -> str:
    """Convert PIL image to base64, resizing if needed."""
    # Convert to RGB (handles RGBA, P, etc.)
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85)

    # Resize if too large
    while buf.tell() > max_bytes and img.size[0] > 200:
        img = img.resize(
            (img.size[0] // 2, img.size[1] // 2),
            Image.LANCZOS
        )
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=80)

    return base64.b64encode(buf.getvalue()).decode("utf-8")


def is_meaningful_image(img: Image.Image) -> bool:
    """Skip tiny decorative images like bullets, icons, borders."""
    w, h = img.size
    return w >= MIN_IMAGE_SIZE and h >= MIN_IMAGE_SIZE


# ── PDF extractor ──────────────────────────────────────────────────────────────

def extract_images_from_pdf(file_path: str) -> list[dict]:
    """
    Extract all images from a PDF with page number and nearby text context.
    Returns list of {page, index, image (PIL), context_text}
    """
    results = []
    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc, start=1):
        # Get nearby text for context
        page_text = page.get_text("text")[:500]

        image_list = page.get_images(full=True)
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes  = base_image["image"]
                img        = Image.open(io.BytesIO(img_bytes))

                if is_meaningful_image(img):
                    results.append({
                        "page":         page_num,
                        "index":        img_idx,
                        "image":        img,
                        "context_text": page_text,
                    })
            except Exception as e:
                print(f"  ⚠ Skipping image xref={xref} on page {page_num}: {e}")

    doc.close()
    return results


# ── PPTX extractor ─────────────────────────────────────────────────────────────

def extract_images_from_pptx(file_path: str) -> list[dict]:
    """Extract images from PowerPoint slides with slide number and text context."""
    results = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Collect all text on this slide for context
        slide_text = " ".join(
            shape.text for shape in slide.shapes
            if shape.has_text_frame
        )[:500]

        img_idx = 0
        for shape in slide.shapes:
            # Check for picture shapes
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    img       = Image.open(io.BytesIO(img_bytes))

                    if is_meaningful_image(img):
                        results.append({
                            "page":         slide_num,
                            "index":        img_idx,
                            "image":        img,
                            "context_text": slide_text,
                        })
                        img_idx += 1
                except Exception as e:
                    print(f"  ⚠ Skipping image on slide {slide_num}: {e}")

    return results


# ── DOCX extractor ─────────────────────────────────────────────────────────────

def extract_images_from_docx(file_path: str) -> list[dict]:
    """Extract images embedded in a Word document."""
    results = []
    doc     = Document(file_path)

    # Get full doc text for context (first 1000 chars)
    full_text = " ".join(p.text for p in doc.paragraphs)[:1000]

    img_idx = 0
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                img_bytes = rel.target_part.blob
                img       = Image.open(io.BytesIO(img_bytes))

                if is_meaningful_image(img):
                    results.append({
                        "page":         1,   # DOCX has no page concept pre-render
                        "index":        img_idx,
                        "image":        img,
                        "context_text": full_text,
                    })
                    img_idx += 1
            except Exception as e:
                print(f"  ⚠ Skipping image in docx: {e}")

    return results


# ── Main ingestion function ────────────────────────────────────────────────────

async def ingest_document_images(
    file_path:     str,
    document_id:   str,
    document_name: str,
    uploader_name: str,
    program:       str,
    subject:       str,
) -> list[ImageChunk]:
    """
    Full pipeline: extract images → describe → embed → return ImageChunks.
    Call this during document upload in your FastAPI endpoint.
    """
    path      = Path(file_path)
    extension = path.suffix.lower()
    chunks    = []

    print(f"\n📄 Processing: {document_name} ({extension})")

    # 1. Extract raw images based on file type
    if extension == ".pdf":
        raw_images = extract_images_from_pdf(file_path)
        source_type = "pdf"
    elif extension in (".pptx", ".ppt"):
        raw_images = extract_images_from_pptx(file_path)
        source_type = "pptx"
    elif extension in (".docx", ".doc"):
        raw_images = extract_images_from_docx(file_path)
        source_type = "docx"
    else:
        print(f"  ℹ Unsupported file type: {extension}")
        return []

    print(f"  🖼  Found {len(raw_images)} meaningful images")

    # 2. For each image: describe + embed
    for item in raw_images:
        img: Image.Image = item["image"]
        w, h = img.size

        # Convert to base64
        b64 = pil_to_b64(img)

        print(f"  🔍 Describing image {item['index']+1}/{len(raw_images)} "
              f"(page/slide {item['page']}, {w}×{h}px)...")

        # Vision description
        description = await describe_image(b64, item["context_text"])
        print(f"     → {description[:100]}...")

        # Embed the description (text embedding — works great for RAG)
        embedding_text = (
            f"Image from {document_name} by {uploader_name} "
            f"({program} - {subject}): {description}"
        )
        embedding = await embed_text(embedding_text)

        chunk = ImageChunk(
            chunk_id      = str(uuid.uuid4()),
            document_id   = document_id,
            document_name = document_name,
            uploader_name = uploader_name,
            program       = program,
            subject       = subject,
            source_type   = source_type,
            page_or_slide = item["page"],
            image_index   = item["index"],
            description   = description,
            image_b64     = b64,
            width         = w,
            height        = h,
            embedding     = embedding,
        )
        chunks.append(chunk)

    print(f"  ✅ Generated {len(chunks)} image chunks")
    return chunks


# ── Student image query ────────────────────────────────────────────────────────

async def query_image(
    student_image_b64: str,
    program_filter: Optional[str] = None,
) -> dict:
    """
    When a student uploads an image:
    1. Describe what they uploaded
    2. Return description + embedding for vector search
    """
    print("🔎 Analyzing student-uploaded image...")

    description = await describe_image(
        student_image_b64,
        context_hint="A student is asking about this object. What is it exactly?"
    )

    # Add program context to narrow search
    search_text = description
    if program_filter:
        search_text = f"{program_filter} context: {description}"

    embedding = await embed_text(search_text)

    return {
        "description": description,
        "embedding":   embedding,
        "search_text": search_text,
    }


# ── Supabase storage helper ────────────────────────────────────────────────────

def chunk_to_supabase_row(chunk: ImageChunk) -> dict:
    """
    Convert an ImageChunk to a row ready for Supabase insertion.
    Table: image_chunks
    """
    d = asdict(chunk)
    # embedding stored separately in pgvector column
    embedding = d.pop("embedding")
    return {
        **d,
        "embedding": embedding,  # pgvector accepts list[float] directly
    }